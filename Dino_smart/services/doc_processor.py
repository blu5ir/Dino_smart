import io
import re
import os
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
from PIL import Image

class DocumentProcessor:
    @staticmethod
    def extract_text_from_pdf(file_bytes: bytes) -> str:
        """Extract text from PDF file bytes."""
        pdf_file = io.BytesIO(file_bytes)
        reader = PdfReader(pdf_file)
        text_list = []
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text_list.append(page_text)
        return "\n\n".join(text_list)

    @staticmethod
    def extract_text_from_docx(file_bytes: bytes) -> str:
        """Extract text from DOCX file bytes, including paragraphs and tables."""
        docx_file = io.BytesIO(file_bytes)
        doc = DocxDocument(docx_file)
        text_list = []
        
        # Read paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_list.append(para.text)
                
        # Read tables
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    text_list.append(" | ".join(row_text))
                    
        return "\n\n".join(text_list)

    @staticmethod
    def extract_text_from_pptx(file_bytes: bytes) -> str:
        """Extract text from PPTX slides."""
        pptx_file = io.BytesIO(file_bytes)
        prs = PptxPresentation(pptx_file)
        text_list = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            slide_text.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            text_list.append("\n".join(slide_text))
            
        return "\n\n".join(text_list)

    @staticmethod
    def extract_text_from_image(file_bytes: bytes, file_name: str, provider) -> str:
        """Perform OCR on an image using the active AI provider's vision system."""
        ext = os.path.splitext(file_name)[1].lower().replace(".", "")
        mime_type = f"image/{ext}"
        if ext not in ["png", "jpg", "jpeg", "webp"]:
            mime_type = "image/png"  # Default fallback
            
        prompt = (
            "You are an expert OCR engine. Extract all readable text from this image exactly as it is shown. "
            "Retain headings, paragraphs, bullet points, and equations. "
            "Do not add explanations, notes, metadata or summaries. Output ONLY the extracted text."
        )
        
        ocr_result = provider.generate_image_ocr(prompt, file_bytes, mime_type)
        return ocr_result

    @staticmethod
    def clean_text(text: str) -> str:
        """Remove duplicate text headers, redundant white-spaces, and clean formatting."""
        # Split into lines
        lines = text.split("\n")
        cleaned_lines = []
        seen_lines = set()
        
        for line in lines:
            trimmed = line.strip()
            if not trimmed:
                cleaned_lines.append("")
                continue
            
            # Simple duplicate suppression for slide templates, footers, page numbers
            # If line is short and we've seen it, skip it
            if len(trimmed) < 60:
                if trimmed.lower() in seen_lines:
                    continue
                seen_lines.add(trimmed.lower())
            
            cleaned_lines.append(line)
            
        # Reconstruct and normalize whitespace
        joined = "\n".join(cleaned_lines)
        # Replace 3 or more consecutive newlines with 2 newlines
        joined = re.sub(r'\n{3,}', '\n\n', joined)
        return joined

    @staticmethod
    def detect_headings(text: str) -> list:
        """Identify potential major headings/sections from text using simple patterns."""
        headings = []
        lines = text.split("\n")
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Pattern 1: Markdown headings
            if line.startswith("#"):
                headings.append(line.lstrip("#").strip())
                continue
                
            # Pattern 2: Short numbered headings (e.g. 1. Introduction, 1.2 Background, Section I: ...)
            numbered_match = re.match(r'^((?:Section\s+)?[I|V|X|L|C\d]+[\.\:\s\-]+[A-Z][a-zA-Z\s]{3,50})$', line)
            if numbered_match:
                headings.append(numbered_match.group(1))
                continue
                
            # Pattern 3: Short all-caps lines (potential heading)
            if len(line) > 3 and len(line) < 60 and line.isupper() and not line.endswith(".") and re.match(r'^[A-Z\s\-\,\&\:]+$', line):
                headings.append(line)
                continue
                
        # Deduplicate headings list
        return list(dict.fromkeys(headings))[:15] # Return top 15 detected headings

    @classmethod
    def process_document(cls, file_bytes: bytes, file_name: str, file_type: str, provider, db_manager, workspace_id: int) -> dict:
        """
        Extracts, cleans, summarizes, detects topics, and stores the document metadata in SQLite.
        Returns a summary dictionary of what was processed.
        """
        # Extract text based on file format
        raw_text = ""
        ft = file_type.lower().strip()
        
        if ft in ["pdf"]:
            raw_text = cls.extract_text_from_pdf(file_bytes)
        elif ft in ["docx", "doc"]:
            raw_text = cls.extract_text_from_docx(file_bytes)
        elif ft in ["pptx", "ppt"]:
            raw_text = cls.extract_text_from_pptx(file_bytes)
        elif ft in ["txt", "md", "markdown"]:
            raw_text = file_bytes.decode("utf-8", errors="ignore")
        elif ft in ["png", "jpg", "jpeg", "webp"]:
            raw_text = cls.extract_text_from_image(file_bytes, file_name, provider)
        else:
            raw_text = file_bytes.decode("utf-8", errors="ignore")
            
        cleaned_text = cls.clean_text(raw_text)
        # Save raw file to upload storage
        from database.db_manager import WORKSPACE_DIR
        upload_path = os.path.join(WORKSPACE_DIR, "uploads", file_name)
        with open(upload_path, "wb") as f:
            f.write(file_bytes)

        # Call AI for metadata extraction if provider exists
        summary = "No summary generated."
        topics = []
        
        if provider and len(cleaned_text.strip()) > 50:
            # 1. Summary generation
            summary_prompt = (
                "You are an assistant. Provide a brief 2-3 sentence summary of the following document content. "
                "Do not include meta-commentary. Focus purely on the main subject.\n\n"
                f"Document Content:\n{cleaned_text[:6000]}"
            )
            summary = provider.generate_text(summary_prompt, system_instruction="Summarize the provided content.")
            
            # 2. Topic extraction
            topics_prompt = (
                "Review the following document and identify 3 to 6 major topic terms or subjects. "
                "Return them as a simple comma-separated list. Do not write any introduction or conclusion.\n\n"
                f"Document Content:\n{cleaned_text[:6000]}"
            )
            topics_raw = provider.generate_text(topics_prompt, system_instruction="Extract comma-separated topics.")
            if topics_raw and "error" not in topics_raw.lower():
                topics = [t.strip() for t in topics_raw.split(",") if t.strip()]
                # Clean topics lists
                topics = [t for t in topics if len(t) < 40]
            
        # If topics couldn't be extracted, use headings or default
        if not topics:
            topics = cls.detect_headings(cleaned_text)[:4]
        if not topics:
            topics = ["General Study"]

        # Store in Database
        doc_id = db_manager.add_document(
            workspace_id=workspace_id,
            name=file_name,
            file_type=ft,
            content=cleaned_text,
            summary=summary,
            topics=topics
        )
        
        return {
            "id": doc_id,
            "name": file_name,
            "file_type": ft,
            "text_length": len(cleaned_text),
            "summary": summary,
            "topics": topics
        }
